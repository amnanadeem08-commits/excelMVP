"""Export system.

Produces client deliverables from the analyzed data:
- Cleaned dataset as Excel or CSV.
- A multi-sheet Excel summary report (overview, KPIs, insights, pivots).
- Chart PNG images.
- Branded PDF report and PowerPoint deck.

CUSTOMIZE:
- Rebrand the PDF/PPT cover text and colors in the theme blocks below.
- Add or remove summary sheets in ``export_summary_report``.
"""

from __future__ import annotations

import html
import json
import os
import tempfile
import time
import tracemalloc
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any

import pandas as pd
from openpyxl.drawing.image import Image as XLImage
from PIL import Image as PILImage
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image as RLImage
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


CHART_WIDTH_PX = 1920
CHART_HEIGHT_PX = 1080
CHART_SCALE = 2
BASE_DPI = 96


@dataclass
class ExportStage:
    stage: str
    status: str
    details: str
    ts_ms: int


def _trace(trace: list[ExportStage], stage: str, status: str, details: str) -> None:
    trace.append(ExportStage(stage=stage, status=status, details=details, ts_ms=int(time.time() * 1000)))


def _safe_sheet_name(value: str, fallback: str = "Sheet") -> str:
    name = (value or fallback).replace("/", "-").replace("\\", "-").replace("*", "-").replace("?", "")
    name = name.replace("[", "(").replace("]", ")").replace(":", "-")
    return name[:31] or fallback


def _verify_written_image(path: str) -> tuple[int, int]:
    with PILImage.open(path) as img:
        img.verify()
    with PILImage.open(path) as img:
        width, height = img.size
    return width, height


def _cleanup_temp_images(images: list[dict[str, Any]], trace: list[ExportStage], stage: str) -> None:
    removed = 0
    for image in images:
        path = image.get("path")
        if not path:
            continue
        try:
            if os.path.exists(path):
                os.remove(path)
                removed += 1
        except Exception as exc:
            _trace(trace, stage, "failed", f"Cleanup failed for {path}: {exc}")
    _trace(trace, stage, "ok", f"Temporary files removed: {removed}")


# ---------------------------------------------------------------------------
# Spreadsheet exports
# ---------------------------------------------------------------------------
def dataframe_to_csv(df: pd.DataFrame) -> bytes:
    """Export a DataFrame to CSV bytes."""
    return df.to_csv(index=False).encode("utf-8")


def dataframe_to_excel(
    df: pd.DataFrame,
    pivots: list[Any] | None = None,
    charts: list[Any] | None = None,
    summary: dict[str, Any] | None = None,
    kpis: dict[str, dict[str, str]] | None = None,
    insights: list[str] | None = None,
    recommendations: list[str] | None = None,
    return_metadata: bool = False,
) -> bytes | tuple[bytes, dict[str, Any]]:
    """Export workbook with dashboard visuals, summary, pivots, charts, and raw data."""
    trace: list[ExportStage] = []
    started = time.perf_counter()
    tracemalloc.start()

    _trace(trace, "dataset", "ok", f"Rows={len(df)}, Cols={df.shape[1]}")
    charts = charts or []
    pivots = pivots or []
    chart_images, failures = export_chart_images(charts, trace=trace)

    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        # Required structure.
        df.to_excel(writer, index=False, sheet_name="Raw Data")

        summary_sheet = pd.DataFrame(
            {
                "Metric": ["Rows", "Columns", "Missing %", "Completeness %"],
                "Value": [
                    summary.get("rows", len(df)) if summary else len(df),
                    summary.get("columns", df.shape[1]) if summary else df.shape[1],
                    summary.get("missing_pct", "n/a") if summary else "n/a",
                    round(100 - float(summary.get("missing_pct", 0)), 1) if summary and isinstance(summary.get("missing_pct", 0), (int, float)) else "n/a",
                ],
            }
        )
        summary_sheet.to_excel(writer, index=False, sheet_name="Summary")

        dashboard_rows = []
        for label, metric in (kpis or {}).items():
            dashboard_rows.append({"KPI": label, "Value": metric.get("value", ""), "Note": metric.get("note", "")})
        if insights:
            for idx, text in enumerate(insights[:6], start=1):
                dashboard_rows.append({"KPI": f"Insight {idx}", "Value": text, "Note": ""})
        if recommendations:
            for idx, text in enumerate(recommendations[:6], start=1):
                dashboard_rows.append({"KPI": f"Recommendation {idx}", "Value": text, "Note": ""})
        pd.DataFrame(dashboard_rows or [{"KPI": "No KPI", "Value": "n/a", "Note": ""}]).to_excel(
            writer, index=False, sheet_name="Dashboard"
        )

        pd.DataFrame(
            [{"Pivot": getattr(p, "title", f"Pivot {i + 1}"), "Rows": len(getattr(p, "data", []))} for i, p in enumerate(pivots)]
            or [{"Pivot": "No pivots", "Rows": 0}]
        ).to_excel(writer, index=False, sheet_name="Pivot Tables")

        pd.DataFrame(
            [{"Chart": img["title"], "Width(px)": img["width_px"], "Height(px)": img["height_px"], "DPI": img["dpi"]} for img in chart_images]
            or [{"Chart": "No charts", "Width(px)": 0, "Height(px)": 0, "DPI": 0}]
        ).to_excel(writer, index=False, sheet_name="Charts")

        for i, pivot in enumerate(pivots, start=1):
            sheet_name = _safe_sheet_name(f"Pivot_{i}_{getattr(pivot, 'title', i)}", fallback=f"Pivot_{i}")
            getattr(pivot, "data", pd.DataFrame()).to_excel(writer, index=False, sheet_name=sheet_name)

        book = writer.book
        dashboard_ws = book["Dashboard"]
        charts_ws = book["Charts"]

        embed_count = 0
        if chart_images:
            # Dashboard sheet: embed lead visuals.
            for idx, img in enumerate(chart_images[:2]):
                xl_img = XLImage(img["path"])
                xl_img.width = 860
                xl_img.height = 480
                dashboard_ws.add_image(xl_img, f"E{2 + idx * 26}")
                embed_count += 1

            # Charts sheet: gallery of all visuals.
            row_anchor = 8
            for img in chart_images:
                xl_img = XLImage(img["path"])
                xl_img.width = 960
                xl_img.height = 540
                charts_ws.add_image(xl_img, f"A{row_anchor}")
                row_anchor += 30
                embed_count += 1

        _trace(trace, "excel_export", "ok", f"Workbook sheets={len(book.sheetnames)}, chart_embeds={embed_count}")

    excel_bytes = output.getvalue()
    _cleanup_temp_images(chart_images, trace, "excel_cleanup")
    elapsed_ms = int((time.perf_counter() - started) * 1000)
    _, peak_mem = tracemalloc.get_traced_memory()
    tracemalloc.stop()

    metadata = {
        "pipeline": "excel",
        "chart_count_input": len(charts),
        "chart_count_rendered": len(chart_images),
        "chart_failures": failures,
        "chart_embedded": embed_count,
        "trace": [stage.__dict__ for stage in trace],
        "elapsed_ms": elapsed_ms,
        "peak_memory_mb": round(peak_mem / (1024 * 1024), 2),
    }
    if return_metadata:
        return excel_bytes, metadata
    return excel_bytes


def export_summary_report(
    summary: dict,
    cleaning_report: dict,
    kpis: dict[str, dict[str, str]],
    insights: list[str],
    recommendations: list[str],
    pivots: list[Any] | None = None,
) -> bytes:
    """Export a structured, multi-sheet Excel summary report.

    Sheets: Overview, KPIs, Cleaning Report, Insights, Recommendations, and
    one sheet per pivot table.
    """
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        overview = pd.DataFrame(
            {
                "Metric": ["Rows", "Columns", "Missing %", "Completeness %", "Numeric Fields", "Categorical Fields", "Datetime Fields"],
                "Value": [
                    summary["rows"],
                    summary["columns"],
                    summary["missing_pct"],
                    round(100 - summary["missing_pct"], 1),
                    summary["numeric_count"],
                    summary["categorical_count"],
                    summary["datetime_count"],
                ],
            }
        )
        overview.to_excel(writer, index=False, sheet_name="Overview")

        kpi_df = pd.DataFrame(
            [{"KPI": label, "Value": metric["value"], "Detail": metric.get("note", "")} for label, metric in kpis.items()]
        )
        kpi_df.to_excel(writer, index=False, sheet_name="KPIs")

        pd.DataFrame([cleaning_report]).T.reset_index().rename(
            columns={"index": "Step", 0: "Value"}
        ).to_excel(writer, index=False, sheet_name="Cleaning Report")

        pd.DataFrame({"Insight": insights}).to_excel(writer, index=False, sheet_name="Insights")
        pd.DataFrame({"Recommendation": recommendations}).to_excel(writer, index=False, sheet_name="Recommendations")

        for pivot in pivots or []:
            sheet_name = pivot.title[:31].replace("/", "-").replace("\\", "-")
            pivot.data.to_excel(writer, index=False, sheet_name=sheet_name)

    return output.getvalue()


# ---------------------------------------------------------------------------
# Chart images
# ---------------------------------------------------------------------------
def export_chart_images(charts: list[Any], trace: list[ExportStage] | None = None) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    """Render charts to high-resolution PNG/SVG and verify readability.

    Returns ``(exported_images, failures)``. Failures are never swallowed.
    """
    trace = trace if trace is not None else []
    exported: list[dict[str, Any]] = []
    failures: list[dict[str, Any]] = []

    _trace(trace, "chart_object_creation", "ok", f"Charts received: {len(charts)}")

    for idx, chart in enumerate(charts):
        title = getattr(chart, "title", f"Chart {idx + 1}")
        figure = getattr(chart, "figure", None)
        if figure is None:
            msg = "Missing figure object"
            failures.append({"index": idx, "title": title, "stage": "object", "error": msg})
            _trace(trace, "chart_object_verify", "failed", f"{title}: {msg}")
            continue

        try:
            figure.to_dict()
            _trace(trace, "chart_render_verify", "ok", f"{title}: figure object is serializable")
        except Exception as exc:
            failures.append({"index": idx, "title": title, "stage": "render", "error": str(exc)})
            _trace(trace, "chart_render_verify", "failed", f"{title}: {exc}")
            continue

        png_bytes: bytes | None = None
        svg_bytes: bytes | None = None
        png_path = ""
        width = 0
        height = 0
        dpi = BASE_DPI * CHART_SCALE

        try:
            png_bytes = figure.to_image(
                format="png",
                width=CHART_WIDTH_PX,
                height=CHART_HEIGHT_PX,
                scale=CHART_SCALE,
            )
            temp_png = tempfile.NamedTemporaryFile(prefix="excel_ai_chart_", suffix=".png", delete=False)
            temp_png.write(png_bytes)
            temp_png.close()
            png_path = temp_png.name
            width, height = _verify_written_image(png_path)
            _trace(trace, "chart_png_verify", "ok", f"{title}: {png_path} ({width}x{height}px @ {dpi}dpi)")
        except Exception as exc:
            failures.append({"index": idx, "title": title, "stage": "png", "error": str(exc)})
            _trace(trace, "chart_png_verify", "failed", f"{title}: {exc}")
            continue

        try:
            svg_bytes = figure.to_image(
                format="svg",
                width=CHART_WIDTH_PX,
                height=CHART_HEIGHT_PX,
                scale=1,
            )
            _trace(trace, "chart_svg_verify", "ok", f"{title}: SVG bytes={len(svg_bytes)}")
        except Exception as exc:
            failures.append({"index": idx, "title": title, "stage": "svg", "error": str(exc)})
            _trace(trace, "chart_svg_verify", "failed", f"{title}: {exc}")

        exported.append(
            {
                "title": title,
                "bytes": png_bytes,
                "svg_bytes": svg_bytes,
                "path": png_path,
                "width_px": width,
                "height_px": height,
                "dpi": dpi,
                "readable": bool(width and height and png_path),
            }
        )

    _trace(trace, "chart_image_generation", "ok", f"Exported={len(exported)}, Failed={len(failures)}")
    return exported, failures


# ---------------------------------------------------------------------------
# PDF report
# ---------------------------------------------------------------------------
def make_custom_pdf_report(
    summary: dict,
    cleaning_report: dict,
    insights: list[str],
    recommendations: list[str],
    pivots: list[Any],
    charts: list[Any],
    include_tables: bool,
    include_charts: bool,
    return_metadata: bool = False,
) -> bytes | tuple[bytes, dict[str, Any]]:
    """Build a branded, client-ready PDF report."""
    trace: list[ExportStage] = []
    started = time.perf_counter()
    tracemalloc.start()
    _trace(trace, "pdf_export_start", "ok", f"include_charts={include_charts}, include_tables={include_tables}")

    chart_images, failures = export_chart_images(charts, trace=trace) if include_charts else ([], [])
    if include_charts and not chart_images:
        _trace(trace, "pdf_chart_verify", "failed", "No chart images were generated for PDF export")
        raise RuntimeError("PDF export aborted: no chart images generated; check chart pipeline errors.")

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=28, leftMargin=28, topMargin=28, bottomMargin=28)
    styles = getSampleStyleSheet()
    styles["Title"].fontSize = 24
    styles["Title"].textColor = colors.HexColor("#111827")
    styles["Heading1"].textColor = colors.HexColor("#111827")
    styles["Heading2"].textColor = colors.HexColor("#1F2937")
    styles["BodyText"].fontSize = 9.5
    styles["BodyText"].leading = 13
    story = [
        _pdf_cover_block("Excel AI Analytics Assistant", "Client-ready business intelligence report"),
        Spacer(1, 12),
        _pdf_kpi_cards(summary),
        Spacer(1, 14),
    ]

    if include_charts and chart_images:
        story.append(Paragraph("Dashboard Visual Snapshot", styles["Heading1"]))
        story.append(Spacer(1, 6))
        for idx, chart_image in enumerate(chart_images[:2]):
            story.append(Paragraph(html.escape(chart_image["title"]), styles["Heading2"]))
            story.append(RLImage(BytesIO(chart_image["bytes"]), width=7.2 * inch, height=4.05 * inch))
            if insights:
                story.append(Paragraph(f"Insight: {html.escape(insights[idx % len(insights)])}", styles["BodyText"]))
            story.append(Spacer(1, 6))

    story.append(Paragraph("Executive Summary", styles["Heading1"]))

    for insight in insights[:5]:
        story.append(Paragraph(f"- {html.escape(insight)}", styles["BodyText"]))
        story.append(Spacer(1, 4))

    story.extend([
        Spacer(1, 10),
        Paragraph("Cleaning Report", styles["Heading2"]),
        _styled_table([[key.replace("_", " ").title(), str(value)] for key, value in cleaning_report.items()]),
        Spacer(1, 12),
        Paragraph("Recommendations", styles["Heading2"]),
    ])
    for recommendation in recommendations[:8]:
        story.append(Paragraph(f"- {html.escape(recommendation)}", styles["BodyText"]))
        story.append(Spacer(1, 4))

    if include_tables and pivots:
        story.append(PageBreak())
        story.append(Paragraph("Pivot Tables", styles["Heading1"]))
        for pivot in pivots[:6]:
            story.append(Spacer(1, 8))
            story.append(Paragraph(html.escape(pivot.title), styles["Heading2"]))
            story.append(_dataframe_pdf_table(pivot.data.head(12)))

    if include_charts and chart_images:
        story.append(PageBreak())
        story.append(Paragraph("Visual Analytics", styles["Heading1"]))
        for idx, chart_image in enumerate(chart_images):
            if idx:
                story.append(PageBreak())
            story.append(Spacer(1, 8))
            story.append(Paragraph(html.escape(chart_image["title"]), styles["Heading2"]))
            story.append(RLImage(BytesIO(chart_image["bytes"]), width=7.2 * inch, height=4.3 * inch))
            if insights:
                story.append(Spacer(1, 4))
                story.append(Paragraph(f"Insight: {html.escape(insights[idx % len(insights)])}", styles["BodyText"]))

    doc.build(story)
    pdf_bytes = buffer.getvalue()
    _cleanup_temp_images(chart_images, trace, "pdf_cleanup")
    elapsed_ms = int((time.perf_counter() - started) * 1000)
    _, peak_mem = tracemalloc.get_traced_memory()
    tracemalloc.stop()
    _trace(trace, "pdf_export_complete", "ok", f"bytes={len(pdf_bytes)}, charts={len(chart_images)}")

    metadata = {
        "pipeline": "pdf",
        "chart_count_input": len(charts),
        "chart_count_rendered": len(chart_images),
        "chart_failures": failures,
        "chart_embedded": len(chart_images),
        "trace": [stage.__dict__ for stage in trace],
        "elapsed_ms": elapsed_ms,
        "peak_memory_mb": round(peak_mem / (1024 * 1024), 2),
    }
    if return_metadata:
        return pdf_bytes, metadata
    return pdf_bytes


def _pdf_cover_block(title: str, subtitle: str) -> Table:
    table = Table(
        [[Paragraph(
            f"<font color='white' size='18'><b>{html.escape(title)}</b></font><br/>"
            f"<font color='#CBD5E1' size='10'>{html.escape(subtitle)}</font>",
            getSampleStyleSheet()["BodyText"],
        )]],
        colWidths=[7.25 * inch],
        rowHeights=[0.85 * inch],
    )
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#111827")),
                ("TEXTCOLOR", (0, 0), (-1, -1), colors.white),
                ("LEFTPADDING", (0, 0), (-1, -1), 14),
                ("RIGHTPADDING", (0, 0), (-1, -1), 14),
                ("TOPPADDING", (0, 0), (-1, -1), 13),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 13),
            ]
        )
    )
    return table


def _pdf_kpi_cards(summary: dict) -> Table:
    data = [[
        _kpi_card_para("Records", f"{summary['rows']:,}"),
        _kpi_card_para("Columns", f"{summary['columns']:,}"),
        _kpi_card_para("Completeness", f"{100 - summary['missing_pct']:.1f}%"),
        _kpi_card_para("Measures", f"{summary['numeric_count']:,}"),
    ]]
    table = Table(data, colWidths=[1.72 * inch] * 4, rowHeights=[0.72 * inch], hAlign="LEFT")
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F8FAFC")),
                ("BOX", (0, 0), (-1, -1), 0.35, colors.HexColor("#D8E1F0")),
                ("INNERGRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#E5EAF3")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("PADDING", (0, 0), (-1, -1), 8),
            ]
        )
    )
    return table


def _kpi_card_para(label: str, value: str) -> Paragraph:
    return Paragraph(
        f"<font size='7' color='#64748B'><b>{html.escape(label.upper())}</b></font><br/>"
        f"<font size='16' color='#111827'><b>{html.escape(value)}</b></font>",
        getSampleStyleSheet()["BodyText"],
    )


def _dataframe_pdf_table(df: pd.DataFrame) -> Table:
    limited = df.head(12).iloc[:, :6].copy()
    rows = [[str(col)[:24] for col in limited.columns]]
    for _, row in limited.iterrows():
        rows.append([str(value)[:34] for value in row])
    table = Table(rows, repeatRows=1)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EEF4FF")),
                ("TEXTCOLOR", (0, 0), (-1, -1), colors.HexColor("#172033")),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#D0D5DD")),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 7),
                ("PADDING", (0, 0), (-1, -1), 4),
            ]
        )
    )
    return table


def _styled_table(rows: list[list[str]]) -> Table:
    table = Table(rows, hAlign="LEFT", colWidths=[180, 300])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EEF4FF")),
                ("TEXTCOLOR", (0, 0), (-1, -1), colors.HexColor("#172033")),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#D0D5DD")),
                ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
                ("PADDING", (0, 0), (-1, -1), 7),
            ]
        )
    )
    return table


# ---------------------------------------------------------------------------
# PowerPoint deck
# ---------------------------------------------------------------------------
def make_ppt_report(
    summary: dict,
    cleaning_report: dict,
    insights: list[str],
    recommendations: list[str],
    pivots: list[Any],
    charts: list[Any],
    include_tables: bool,
    include_charts: bool,
    return_metadata: bool = False,
) -> bytes | tuple[bytes, dict[str, Any]]:
    """Build a branded, client-ready PowerPoint deck."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    trace: list[ExportStage] = []
    started = time.perf_counter()
    tracemalloc.start()
    chart_images, failures = export_chart_images(charts, trace=trace) if include_charts else ([], [])
    if include_charts and not chart_images:
        _trace(trace, "ppt_chart_verify", "failed", "No chart images were generated for PPT export")
        raise RuntimeError("PPT export aborted: no chart images generated; check chart pipeline errors.")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    theme = {
        "ink": RGBColor(17, 24, 39),
        "muted": RGBColor(100, 116, 139),
        "blue": RGBColor(37, 99, 235),
        "cyan": RGBColor(8, 145, 178),
        "green": RGBColor(5, 150, 105),
        "amber": RGBColor(217, 119, 6),
        "panel": RGBColor(248, 250, 252),
        "line": RGBColor(226, 232, 240),
        "navy": RGBColor(15, 23, 42),
    }

    def set_bg(slide, color=RGBColor(255, 255, 255)):
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = color

    def add_title(slide, title: str, subtitle: str = ""):
        set_bg(slide)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme["blue"]
        bar.line.fill.background()
        box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = theme["ink"]
        if subtitle:
            sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.92), Inches(11.8), Inches(0.35))
            sp = sub.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(11)
            sp.font.color.rgb = theme["muted"]

    def add_card(slide, x, y, w, h, label: str, value: str, accent):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = theme["line"]
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()
        label_box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.16), w - Inches(0.35), Inches(0.22))
        lp = label_box.text_frame.paragraphs[0]
        lp.text = label.upper()
        lp.font.size = Pt(8)
        lp.font.bold = True
        lp.font.color.rgb = theme["muted"]
        value_box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.45), w - Inches(0.35), Inches(0.42))
        vp = value_box.text_frame.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(18)
        vp.font.bold = True
        vp.font.color.rgb = theme["ink"]

    def add_bullets(slide, x, y, w, h, title: str, items: list[str]):
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(248, 250, 252)
        panel.line.color.rgb = theme["line"]
        title_box = slide.shapes.add_textbox(x, y, w, Inches(0.35))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = title
        tp.font.bold = True
        tp.font.size = Pt(15)
        tp.font.color.rgb = theme["ink"]
        title_box.left = x + Inches(0.25)
        title_box.top = y + Inches(0.18)
        title_box.width = w - Inches(0.5)
        body = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.65), w - Inches(0.5), h - Inches(0.8))
        tf = body.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(items[:7]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(45, 55, 72)

    def add_df_table(slide, df: pd.DataFrame, x, y, w, h):
        table_df = df.head(8).iloc[:, :5].copy()
        rows, cols = table_df.shape[0] + 1, max(table_df.shape[1], 1)
        table = slide.shapes.add_table(rows, cols, x, y, w, h).table
        for col_idx, col in enumerate(table_df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col)[:28]
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(239, 246, 255)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = theme["ink"]
        for row_idx, (_, row) in enumerate(table_df.iterrows(), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)[:40]
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(7)
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.font.color.rgb = RGBColor(51, 65, 85)

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(title_slide, theme["navy"])
    accent = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = theme["navy"]
    accent.line.fill.background()
    ribbon = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = theme["blue"]
    ribbon.line.fill.background()
    title_box = title_slide.shapes.add_textbox(Inches(0.85), Inches(0.85), Inches(8.8), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Excel AI Analytics Assistant"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    sub_box = title_slide.shapes.add_textbox(Inches(0.9), Inches(1.95), Inches(7.5), Inches(0.5))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "Client-ready analytics export with automated insights, pivots, and dashboard visuals"
    sp.font.size = Pt(14)
    sp.font.color.rgb = RGBColor(203, 213, 225)
    kpi_values = [
        ("Records", f"{summary['rows']:,}", theme["blue"]),
        ("Columns", f"{summary['columns']:,}", theme["cyan"]),
        ("Completeness", f"{100 - summary['missing_pct']:.1f}%", theme["green"]),
        ("Measures", f"{summary['numeric_count']:,}", theme["amber"]),
    ]
    for idx, (label, value, color) in enumerate(kpi_values):
        add_card(title_slide, Inches(0.9 + idx * 3.05), Inches(3.05), Inches(2.65), Inches(1.15), label, value, color)
    add_bullets(
        title_slide,
        Inches(0.9),
        Inches(4.75),
        Inches(5.85),
        Inches(1.75),
        "Executive highlights",
        [
            f"{summary['rows']:,} records analyzed across {summary['columns']:,} columns.",
            f"Data completeness score: {100 - summary['missing_pct']:.1f}%.",
            f"{summary['numeric_count']:,} numeric fields and {summary['categorical_count']:,} segment fields detected.",
        ],
    )
    add_bullets(title_slide, Inches(7.05), Inches(4.75), Inches(5.45), Inches(1.75), "AI recommendations", recommendations)

    insight_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(insight_slide, "AI Insights", "Business-friendly findings generated from the uploaded dataset")
    for idx, (label, value, color) in enumerate(kpi_values):
        add_card(insight_slide, Inches(0.65 + idx * 3.05), Inches(1.35), Inches(2.65), Inches(1.0), label, value, color)
    add_bullets(insight_slide, Inches(0.75), Inches(2.7), Inches(5.75), Inches(3.95), "Detected insights", insights)
    add_bullets(insight_slide, Inches(6.75), Inches(2.7), Inches(5.75), Inches(3.95), "Recommended actions", recommendations)

    if include_tables:
        for pivot in pivots:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_title(slide, pivot.title, "Top rows from the automatically generated pivot table")
            add_df_table(slide, pivot.data, Inches(0.65), Inches(1.35), Inches(12.1), Inches(5.55))

    if include_charts:
        for idx, chart_image in enumerate(chart_images):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_title(slide, chart_image["title"], "Dashboard-matched chart exported at high resolution")
            frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.42), Inches(1.08), Inches(12.5), Inches(5.98))
            frame.fill.solid()
            frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
            frame.line.color.rgb = theme["line"]
            # Chart area occupies ~70% of slide surface.
            slide.shapes.add_picture(BytesIO(chart_image["bytes"]), Inches(0.7), Inches(1.4), width=Inches(9.45), height=Inches(5.0))

            insight_box = slide.shapes.add_textbox(Inches(10.35), Inches(1.55), Inches(2.55), Inches(1.7))
            ip = insight_box.text_frame.paragraphs[0]
            ip.text = "Key insight"
            ip.font.size = Pt(12)
            ip.font.bold = True
            ip.font.color.rgb = theme["ink"]
            insight_detail = insight_box.text_frame.add_paragraph()
            insight_detail.text = (insights[idx % len(insights)] if insights else "Insight not available")[:220]
            insight_detail.font.size = Pt(10)
            insight_detail.font.color.rgb = theme["muted"]

            if include_tables and pivots:
                supporting = pivots[idx % len(pivots)].data
                add_df_table(slide, supporting, Inches(10.3), Inches(3.35), Inches(2.6), Inches(2.2))

    output = BytesIO()
    prs.save(output)
    ppt_bytes = output.getvalue()
    _cleanup_temp_images(chart_images, trace, "ppt_cleanup")
    elapsed_ms = int((time.perf_counter() - started) * 1000)
    _, peak_mem = tracemalloc.get_traced_memory()
    tracemalloc.stop()
    _trace(trace, "ppt_export_complete", "ok", f"bytes={len(ppt_bytes)}, charts={len(chart_images)}")

    metadata = {
        "pipeline": "ppt",
        "chart_count_input": len(charts),
        "chart_count_rendered": len(chart_images),
        "chart_failures": failures,
        "chart_embedded": len(chart_images),
        "trace": [stage.__dict__ for stage in trace],
        "elapsed_ms": elapsed_ms,
        "peak_memory_mb": round(peak_mem / (1024 * 1024), 2),
    }
    if return_metadata:
        return ppt_bytes, metadata
    return ppt_bytes


def generate_export_validation_report(
    df: pd.DataFrame,
    summary: dict,
    cleaning_report: dict,
    kpis: dict[str, dict[str, str]],
    insights: list[str],
    recommendations: list[str],
    pivots: list[Any],
    charts: list[Any],
    include_tables: bool,
    include_charts: bool,
) -> str:
    """Run full export validation and return a JSON debug report."""
    trace: list[ExportStage] = []
    started = time.perf_counter()
    tracemalloc.start()

    _trace(trace, "dataset", "ok", f"rows={len(df)}, cols={df.shape[1]}")
    _trace(trace, "cleaning", "ok", f"missing_fixed={cleaning_report.get('missing_values_fixed', 'n/a')}")
    _trace(trace, "pivot_generation", "ok", f"pivot_count={len(pivots)}")

    chart_images, chart_failures = export_chart_images(charts, trace=trace)
    chart_count = len(charts)
    rendered_count = len(chart_images)

    excel_ok = ppt_ok = pdf_ok = False
    excel_meta: dict[str, Any] = {}
    ppt_meta: dict[str, Any] = {}
    pdf_meta: dict[str, Any] = {}
    export_failures: list[str] = []

    try:
        excel_bytes, excel_meta = dataframe_to_excel(
            df,
            pivots=pivots,
            charts=charts,
            summary=summary,
            kpis=kpis,
            insights=insights,
            recommendations=recommendations,
            return_metadata=True,
        )
        excel_ok = len(excel_bytes) > 1000 and excel_meta.get("chart_embedded", 0) > 0
    except Exception as exc:
        export_failures.append(f"excel: {exc}")

    try:
        ppt_bytes, ppt_meta = make_ppt_report(
            summary,
            cleaning_report,
            insights,
            recommendations,
            pivots,
            charts,
            include_tables,
            include_charts,
            return_metadata=True,
        )
        ppt_ok = len(ppt_bytes) > 1000 and (ppt_meta.get("chart_embedded", 0) > 0 if include_charts else True)
    except Exception as exc:
        export_failures.append(f"ppt: {exc}")

    try:
        pdf_bytes, pdf_meta = make_custom_pdf_report(
            summary,
            cleaning_report,
            insights,
            recommendations,
            pivots,
            charts,
            include_tables,
            include_charts,
            return_metadata=True,
        )
        pdf_ok = len(pdf_bytes) > 1000 and (pdf_meta.get("chart_embedded", 0) > 0 if include_charts else True)
    except Exception as exc:
        export_failures.append(f"pdf: {exc}")

    elapsed_ms = int((time.perf_counter() - started) * 1000)
    _, peak_mem = tracemalloc.get_traced_memory()
    tracemalloc.stop()

    report = {
        "status": "ok" if not export_failures else "failed",
        "chart_pipeline": {
            "charts_detected": chart_count,
            "charts_rendered": rendered_count,
            "charts_exported": rendered_count,
            "missing_charts": max(chart_count - rendered_count, 0),
            "failures": chart_failures,
            "chart_details": [
                {
                    "title": img["title"],
                    "path": img["path"],
                    "width_px": img["width_px"],
                    "height_px": img["height_px"],
                    "dpi": img["dpi"],
                    "readable": img["readable"],
                    "svg_supported": bool(img.get("svg_bytes")),
                }
                for img in chart_images
            ],
        },
        "exports": {
            "excel": {
                "ok": excel_ok,
                "charts_embedded": excel_meta.get("chart_embedded", 0),
                "chart_failures": excel_meta.get("chart_failures", []),
            },
            "powerpoint": {
                "ok": ppt_ok,
                "charts_embedded": ppt_meta.get("chart_embedded", 0),
                "chart_failures": ppt_meta.get("chart_failures", []),
            },
            "pdf": {
                "ok": pdf_ok,
                "charts_embedded": pdf_meta.get("chart_embedded", 0),
                "chart_failures": pdf_meta.get("chart_failures", []),
            },
            "storyboard": {
                "ok": ppt_ok,
                "note": "Storyboard uses PowerPoint visual pipeline.",
            },
        },
        "validation_checks": {
            "dashboard_generated": chart_count > 0,
            "charts_generated": rendered_count > 0,
            "chart_files_created": rendered_count > 0,
            "charts_embedded_excel": excel_meta.get("chart_embedded", 0) > 0,
            "charts_embedded_powerpoint": ppt_meta.get("chart_embedded", 0) > 0 if include_charts else True,
            "charts_embedded_pdf": pdf_meta.get("chart_embedded", 0) > 0 if include_charts else True,
            "no_placeholder_tables_replacing_charts": rendered_count > 0 if include_charts else True,
            "no_broken_exports": not export_failures,
        },
        "failures": export_failures,
        "execution": {
            "elapsed_ms": elapsed_ms,
            "peak_memory_mb": round(peak_mem / (1024 * 1024), 2),
        },
        "stage_trace": [stage.__dict__ for stage in trace],
    }

    # Temporary files are no longer needed after report assembly.
    _cleanup_temp_images(chart_images, trace, "report_cleanup")
    return json.dumps(report, indent=2)
