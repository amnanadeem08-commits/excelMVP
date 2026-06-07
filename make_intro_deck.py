"""Generate a LinkedIn-ready portfolio PowerPoint for the project.

Builds a shareable deck for GitHub, LinkedIn posts, and freelance portfolios.
Includes the live repo link and a ready-to-copy LinkedIn caption.

Run:  py make_intro_deck.py
Output:  Excel_MVP_LinkedIn_Portfolio.pptx
"""

from __future__ import annotations

GITHUB_URL = "https://github.com/amnanadeem08-commits/excelMVP"
LINKEDIN_HASHTAGS = "#Python #DataAnalytics #ExcelAutomation #Streamlit #AI #Freelance #OpenSource"

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---- Theme ----------------------------------------------------------------
NAVY = RGBColor(15, 23, 42)
INK = RGBColor(17, 24, 39)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(8, 145, 178)
GREEN = RGBColor(5, 150, 105)
AMBER = RGBColor(217, 119, 6)
PURPLE = RGBColor(124, 58, 237)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(248, 250, 252)
LINE = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)
SLATE = RGBColor(51, 65, 85)

ACCENTS = [BLUE, CYAN, GREEN, AMBER, PURPLE]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---- Helpers --------------------------------------------------------------
def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def rrect(slide, x, y, w, h, color, line=LINE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    """runs: list of (string, size, bold, color) or list of such lists per paragraph."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if runs and not isinstance(runs[0], list):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        string, size, bold, color = para
        run = p.add_run()
        run.text = string
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, size=14, color=SLATE, bullet="•", space=8):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def header(slide, kicker, title):
    set_bg(slide, WHITE)
    rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.12), BLUE)
    text(slide, Inches(0.6), Inches(0.42), Inches(12), Inches(0.3), (kicker.upper(), 12, True, BLUE))
    text(slide, Inches(0.6), Inches(0.72), Inches(12), Inches(0.7), (title, 30, True, INK))
    rect(slide, Inches(0.62), Inches(1.42), Inches(1.1), Inches(0.05), AMBER)


def card(slide, x, y, w, h, accent, title, body, icon=""):
    rrect(slide, x, y, w, h, WHITE, LINE)
    rect(slide, x, y, Inches(0.09), h, accent)
    head = f"{icon}  {title}" if icon else title
    text(slide, x + Inches(0.28), y + Inches(0.18), w - Inches(0.45), Inches(0.4), (head, 15, True, INK))
    text(slide, x + Inches(0.28), y + Inches(0.66), w - Inches(0.45), h - Inches(0.8), (body, 11.5, False, SLATE))


# ===========================================================================
# Slide 1 — Title
# ===========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)
rect(s, Inches(0), Inches(0), Inches(0.18), prs.slide_height, BLUE)
text(s, Inches(0.9), Inches(1.1), Inches(11), Inches(0.4), ("PYTHON · STREAMLIT · PANDAS · AI", 13, True, RGBColor(125, 211, 252)))
text(s, Inches(0.85), Inches(1.7), Inches(11.6), Inches(1.8),
     [["Excel Automation &", 46, True, WHITE], ["AI Reporting Tool", 46, True, WHITE]])
text(s, Inches(0.9), Inches(3.5), Inches(10.8), Inches(1.0),
     ("Upload a spreadsheet and automatically clean data, build dashboards, "
      "generate pivot analytics, and produce client-ready AI business reports.", 17, False, RGBColor(203, 213, 225)))
text(s, Inches(0.92), Inches(4.55), Inches(10.5), Inches(0.45),
     (GITHUB_URL, 14, True, RGBColor(125, 211, 252)))
# bottom KPI ribbon
items = [("6", "Core modules", BLUE), ("4-section", "AI insight report", CYAN), ("PDF · PPT · Excel", "Exports", GREEN), ("Any", "Excel/CSV file", AMBER)]
bx = Inches(0.9)
for val, lbl, acc in items:
    rrect(s, bx, Inches(5.35), Inches(2.75), Inches(1.25), RGBColor(30, 41, 59), None)
    rect(s, bx, Inches(5.35), Inches(0.08), Inches(1.25), acc)
    text(s, bx + Inches(0.25), Inches(5.5), Inches(2.4), Inches(0.5), (val, 22, True, WHITE))
    text(s, bx + Inches(0.25), Inches(6.0), Inches(2.4), Inches(0.4), (lbl, 11, False, RGBColor(148, 163, 184)))
    bx += Inches(2.95)

# ===========================================================================
# Slide 2 — What is it (Intro)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Introduction", "What is this project?")
text(s, Inches(0.62), Inches(1.7), Inches(7.0), Inches(1.4),
     ("A freelance-ready web app that replaces hours of manual Excel work. "
      "It reads any spreadsheet, cleans and structures the data with pandas, "
      "and turns it into interactive dashboards and professional reports — "
      "explained in plain business language by an AI insight layer.", 15, False, SLATE))
rrect(s, Inches(0.62), Inches(3.4), Inches(7.0), Inches(3.4), LIGHT, LINE)
text(s, Inches(0.92), Inches(3.6), Inches(6.4), Inches(0.4), ("Built like a product, not a demo", 15, True, INK))
bullets(s, Inches(0.92), Inches(4.1), Inches(6.4), Inches(2.6), [
    "Simple, client-friendly Streamlit interface",
    "Works with any sales, finance, HR, or inventory file",
    "No manual formulas, pivot tables, or chart building",
    "One-click cleaned data, dashboards, and reports",
    "Optional real AI (OpenAI/Claude) for written narratives",
], size=13)
# right panel cards
card(s, Inches(7.95), Inches(1.7), Inches(4.75), Inches(1.55), BLUE, "The problem", "Small businesses drown in messy Excel files and spend hours building the same reports by hand.")
card(s, Inches(7.95), Inches(3.4), Inches(4.75), Inches(1.55), GREEN, "The solution", "Automate the entire flow: clean → analyze → visualize → explain → export, in seconds.")
card(s, Inches(7.95), Inches(5.1), Inches(4.75), Inches(1.55), AMBER, "The value", "Sellable on Fiverr & Upwork for Excel automation, dashboards, and reporting gigs.")

# ===========================================================================
# Slide 3 — Key Features
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Capabilities", "Key Features")
feats = [
    (BLUE, "Excel Input System", "Upload .xlsx / .xls / .csv, multi-sheet support, automatic structure detection."),
    (CYAN, "Data Cleaning Engine", "Standardize names, remove duplicates, fix types, fill missing values, cap outliers."),
    (GREEN, "Automation Layer", "Pivot-table equivalents via pandas groupby + KPIs (total, average, growth %)."),
    (AMBER, "Interactive Dashboard", "KPI cards, bar/line/pie charts, distributions, and category + date filters."),
    (PURPLE, "AI Insight Layer", "Detects dataset type and writes summary, insights, anomalies, recommendations."),
    (BLUE, "Export System", "Cleaned Excel/CSV, multi-sheet summary, and branded PDF & PowerPoint reports."),
]
x0, y0, w, h, gx, gy = Inches(0.62), Inches(1.75), Inches(3.95), Inches(2.35), Inches(0.18), Inches(0.25)
for i, (acc, t, b) in enumerate(feats):
    r, c = divmod(i, 3)
    card(s, x0 + c * (w + gx), y0 + r * (h + gy), w, h, acc, t, b)

# ===========================================================================
# Slide 4 — How it works (pipeline)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Architecture", "How It Works — The Pipeline")
steps = [
    (BLUE, "1  Upload", "Read Excel/CSV, pick a sheet, detect structure"),
    (CYAN, "2  Clean", "Fix types, dedupe, fill gaps, cap outliers"),
    (GREEN, "3  Analyze", "Groupby pivots, KPIs, growth, correlations"),
    (AMBER, "4  Visualize", "Auto Plotly dashboard with live filters"),
    (PURPLE, "5  Explain", "AI insights: summary, anomalies, actions"),
    (BLUE, "6  Export", "Excel, PDF & PowerPoint deliverables"),
]
bx = Inches(0.62)
bw = Inches(1.95)
for i, (acc, t, b) in enumerate(steps):
    y = Inches(2.55)
    rrect(s, bx, y, bw, Inches(1.95), WHITE, LINE)
    rect(s, bx, y, bw, Inches(0.12), acc)
    text(s, bx + Inches(0.15), y + Inches(0.3), bw - Inches(0.3), Inches(0.5), (t, 15, True, INK))
    text(s, bx + Inches(0.15), y + Inches(0.85), bw - Inches(0.3), Inches(1.0), (b, 10.5, False, SLATE))
    if i < len(steps) - 1:
        text(s, bx + bw - Inches(0.02), y + Inches(0.6), Inches(0.35), Inches(0.5), ("›", 26, True, MUTED), align=PP_ALIGN.CENTER)
    bx += bw + Inches(0.06)
text(s, Inches(0.62), Inches(4.95), Inches(12), Inches(0.4), ("Data flows one direction — every stage is a separate, reusable Python module.", 13, True, MUTED))
# module mapping strip
mods = "data_loader.py  →  data_cleaning.py  →  analytics_engine.py  →  dashboard.py  →  ai_insights.py  →  export_module.py"
rrect(s, Inches(0.62), Inches(5.55), Inches(12.1), Inches(1.1), LIGHT, LINE)
text(s, Inches(0.62), Inches(5.55), Inches(12.1), Inches(1.1), (mods, 13, True, SLATE), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# Slide 5 — AI Insight Layer (the brain)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "The Brain", "AI Insight Layer")
text(s, Inches(0.62), Inches(1.7), Inches(12), Inches(0.5),
     ("generate_ai_insights(df) converts cleaned data into a structured, client-ready report in four sections:", 14, False, SLATE))
ai_cards = [
    (BLUE, "1 · Dataset Summary", "Auto-detects the domain (sales, finance, HR, inventory) and explains what the data represents."),
    (CYAN, "2 · Key Insights", "Top & worst performers, growth/decline trend, and the most important patterns."),
    (AMBER, "3 · Problems / Anomalies", "Unusual spikes & drops vs. trend, abnormal values, and missing-data patterns."),
    (GREEN, "4 · Recommendations", "3–5 actionable suggestions to grow revenue and fix weak areas."),
]
x0, y0, w, h, gx, gy = Inches(0.62), Inches(2.35), Inches(6.0), Inches(1.65), Inches(0.18), Inches(0.2)
for i, (acc, t, b) in enumerate(ai_cards):
    r, c = divmod(i, 2)
    card(s, x0 + c * (w + gx), y0 + r * (h + gy), w, h, acc, t, b)
rrect(s, Inches(0.62), Inches(5.75), Inches(12.1), Inches(1.05), NAVY, None)
text(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(1.05),
     ('Example:  "Sales dropped by 35% in March compared to the average trend."   ·   '
      '"Focus marketing on Product A — it contributes 42% of total revenue."', 13, True, RGBColor(226, 232, 240)),
     anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# Slide 6 — Tech stack & project structure
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Under the Hood", "Tech Stack & Project Structure")
stack = [
    (BLUE, "Streamlit", "Interactive web UI"),
    (CYAN, "Pandas + NumPy", "Cleaning & analytics"),
    (GREEN, "Plotly", "Interactive charts"),
    (AMBER, "ReportLab + python-pptx", "PDF & PPT exports"),
    (PURPLE, "OpenAI / Claude", "Optional LLM insights"),
]
y = Inches(1.8)
for acc, t, b in stack:
    rrect(s, Inches(0.62), y, Inches(5.9), Inches(0.85), WHITE, LINE)
    rect(s, Inches(0.62), y, Inches(0.09), Inches(0.85), acc)
    text(s, Inches(0.9), y + Inches(0.12), Inches(3.0), Inches(0.6), (t, 14, True, INK), anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.5), y, Inches(2.9), Inches(0.85), (b, 12, False, SLATE), anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.98)
# structure panel
rrect(s, Inches(6.85), Inches(1.8), Inches(5.85), Inches(4.9), NAVY, None)
text(s, Inches(7.15), Inches(2.0), Inches(5.4), Inches(0.4), ("Clean, modular codebase", 15, True, WHITE))
struct = [
    "app.py — Streamlit UI orchestrator",
    "data_loader.py — input & sheet detection",
    "data_cleaning.py — pandas cleaning engine",
    "analytics_engine.py — pivots, KPIs, growth",
    "dashboard.py — charts, cards, filters",
    "ai_insights.py — the AI brain",
    "export_module.py — Excel / PDF / PPT",
]
bullets(s, Inches(7.15), Inches(2.55), Inches(5.4), Inches(4.0), struct, size=12.5, color=RGBColor(203, 213, 225), bullet="›", space=9)

# ===========================================================================
# Slide 7 — Use cases / target market
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Who It's For", "Use Cases & Target Market")
uses = [
    (BLUE, "Small Business Reporting", "Monthly sales, revenue, and operations dashboards without an analyst."),
    (CYAN, "Sales Analysis", "Track top products, regions, and trends; spot drops early."),
    (GREEN, "Finance Automation", "Clean transaction exports and produce summary reports fast."),
    (AMBER, "HR Analytics", "Headcount, salary, tenure, and attrition snapshots."),
    (PURPLE, "Freelance Gigs", "Deliverable for Fiverr & Upwork Excel automation jobs."),
    (BLUE, "Recurring Reports", "Reusable analytics pack for weekly/monthly leadership reviews."),
]
x0, y0, w, h, gx, gy = Inches(0.62), Inches(1.75), Inches(3.95), Inches(2.35), Inches(0.18), Inches(0.25)
for i, (acc, t, b) in enumerate(uses):
    r, c = divmod(i, 3)
    card(s, x0 + c * (w + gx), y0 + r * (h + gy), w, h, acc, t, b)

# ===========================================================================
# Slide 8 — LinkedIn portfolio post (copy-paste ready)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "LinkedIn", "Share This Project on LinkedIn")
linkedin_post = (
    "I built an Excel Automation & AI Reporting Tool in Python.\n\n"
    "Upload any Excel/CSV file and the app automatically:\n"
    "- Cleans and structures messy data\n"
    "- Builds interactive dashboards with client brand colors\n"
    "- Generates pivot-table analytics and KPIs\n"
    "- Writes AI-style business insights\n"
    "- Exports Excel, PDF, and PowerPoint reports\n\n"
    f"Open source on GitHub: {GITHUB_URL}\n\n"
    f"{LINKEDIN_HASHTAGS}"
)
rrect(s, Inches(0.62), Inches(1.75), Inches(7.6), Inches(4.9), LIGHT, LINE)
text(s, Inches(0.92), Inches(1.95), Inches(7.0), Inches(0.4), ("Suggested LinkedIn caption", 15, True, INK))
bullets(s, Inches(0.92), Inches(2.45), Inches(7.0), Inches(4.0), [
    "I built an Excel Automation & AI Reporting Tool in Python.",
    "Upload Excel/CSV -> clean data -> dashboards -> AI insights -> exports.",
    "Built for freelance gigs: sales, finance, HR, and small-business reporting.",
    "Features: client branding, custom colors, logo, PDF/PPT deliverables.",
    f"GitHub: {GITHUB_URL}",
    "Stack: Python, Streamlit, Pandas, Plotly, AI insights layer.",
], size=12.5, bullet=">", space=10)
card(s, Inches(8.45), Inches(1.75), Inches(4.25), Inches(2.2), BLUE, "Best for LinkedIn carousel",
     "Export these slides as PDF or images. Post 3-5 slides with the GitHub link in the first comment.")
card(s, Inches(8.45), Inches(4.15), Inches(4.25), Inches(2.5), GREEN, "Hashtags to use",
     LINKEDIN_HASHTAGS)

# ===========================================================================
# Slide 9 — Closing + GitHub CTA
# ===========================================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, NAVY)
rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.16), BLUE)
text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.0), ("From messy spreadsheet to", 30, True, RGBColor(148, 163, 184)))
text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.0), ("client-ready report in seconds.", 38, True, WHITE))
text(s, Inches(0.92), Inches(3.85), Inches(11), Inches(0.6),
     ("Automated cleaning · Branded dashboards · AI insights · Professional exports", 15, False, RGBColor(203, 213, 225)))
rrect(s, Inches(0.9), Inches(4.75), Inches(5.5), Inches(1.15), RGBColor(30, 41, 59), None)
rect(s, Inches(0.9), Inches(4.75), Inches(0.08), Inches(1.15), GREEN)
text(s, Inches(1.2), Inches(4.75), Inches(5.1), Inches(1.15),
     (f"GitHub Repo:\n{GITHUB_URL}", 14, True, WHITE), anchor=MSO_ANCHOR.MIDDLE)
rrect(s, Inches(6.65), Inches(4.75), Inches(5.75), Inches(1.15), RGBColor(30, 41, 59), None)
rect(s, Inches(6.65), Inches(4.75), Inches(0.08), Inches(1.15), CYAN)
text(s, Inches(6.95), Inches(4.75), Inches(5.3), Inches(1.15),
     ("Run locally:\nstreamlit run app.py", 14, True, WHITE), anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.92), Inches(6.2), Inches(11), Inches(0.5),
     ("Star the repo · Share on LinkedIn · Use for Fiverr/Upwork client demos", 13, False, RGBColor(148, 163, 184)))

OUTPUT = "Excel_MVP_LinkedIn_Portfolio.pptx"
LEGACY_OUTPUT = "Excel_Automation_AI_Tool_Presentation.pptx"
prs.save(OUTPUT)
prs.save(LEGACY_OUTPUT)
print(f"Saved {OUTPUT} with {len(prs.slides._sldIdLst)} slides")
print(f"Also saved {LEGACY_OUTPUT}")
